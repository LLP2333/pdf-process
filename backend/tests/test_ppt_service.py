"""ppt_service:PPTX 生成。"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app import ppt_service
from app.schemas import Question, Segment


def test_build_pptx_creates_one_slide_per_question(sample_pdf: Path, tmp_path: Path) -> None:
    out = tmp_path / "out.pptx"
    questions = [
        Question(no=1, segments=[Segment(page=0, y1=120, y2=300)]),
        Question(no=2, segments=[Segment(page=1, y1=120, y2=240)]),
    ]
    made = ppt_service.build_pptx(sample_pdf, out, questions, margin_pt=28.0)
    assert made == 2

    prs = Presentation(out.as_posix())
    assert len(prs.slides) == 2
    # 16:9 标准尺寸
    assert prs.slide_width == Inches(13.333)
    assert prs.slide_height == Inches(7.5)
    # 每张幻灯片至少有一张图片(题区)
    for slide in prs.slides:
        assert any(shape.shape_type == 13 for shape in slide.shapes), "应至少含 1 张图片(shape_type=13)"


def test_build_pptx_skips_empty_questions(sample_pdf: Path, tmp_path: Path) -> None:
    out = tmp_path / "out.pptx"
    questions = [
        Question(no=1, segments=[Segment(page=99, y1=0, y2=10)]),  # 全部越界
        Question(no=2, segments=[Segment(page=0, y1=100, y2=200)]),
    ]
    made = ppt_service.build_pptx(sample_pdf, out, questions, margin_pt=28.0)
    assert made == 1
