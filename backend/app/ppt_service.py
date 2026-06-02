"""按题目-段落方案生成 PPTX:每题一张 16:9 幻灯片,题区图片置顶,下方留白。

PPTX 不直接支持矢量 PDF 嵌入,这里把每段先用 PyMuPDF 高 DPI 渲染成 PNG
再插入,几何与原 PDF 完全一致,清晰度对于讲解投影完全够用。
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

from . import pdf_service
from .schemas import Question

# 16:9 标准尺寸(英寸)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _emu(inches: float) -> Emu:
    return Inches(inches)


def build_pptx(pdf_path: Path, out_path: Path, questions: list[Question], margin_pt: float) -> int:
    """生成一题一张幻灯片的 16:9 PPTX。

    Args:
        pdf_path: 原 PDF(用于裁剪渲染每段)。
        out_path: 目标 PPTX,父目录会自动创建。
        questions: 用户给的切分方案(`no` + 多段)。
        margin_pt: 每张幻灯片四周留白(pt),与 PDF 导出共用同一参数。
            内部按 1in = 72pt 换算为英寸。

    Returns:
        成功生成的幻灯片数(过滤掉空段题目后)。
    """
    prs = Presentation()
    prs.slide_width = _emu(SLIDE_W_IN)
    prs.slide_height = _emu(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]  # 空白版式

    # 把 PDF 用 pt 表示的留白换算到英寸(1in = 72pt)
    margin_in = margin_pt / 72.0
    avail_w_in = SLIDE_W_IN - 2 * margin_in
    avail_h_in = SLIDE_H_IN - 2 * margin_in

    made = 0
    for q in sorted(questions, key=lambda x: x.no):
        pngs = pdf_service.render_segments_to_png(pdf_path, q)
        if not pngs:
            continue

        # 用第一张图的宽度作为整题宽度(各段裁剪框横向同为整页宽)。
        # 但要在 PPTX 中先按可用宽缩放,再按可用高对总高约束二次缩放。
        from PIL import Image  # 局部导入:Pillow 是 python-pptx 的依赖,已可用

        imgs = [Image.open(BytesIO(b)) for b in pngs]
        per_w_px = max(im.width for im in imgs)
        total_h_px = sum(im.height for im in imgs)
        # 像素 → 英寸:按宽度缩放,使图宽 = 可用宽
        scale_in_per_px = avail_w_in / per_w_px
        scaled_h_in = total_h_px * scale_in_per_px
        if scaled_h_in > avail_h_in:
            shrink = avail_h_in / scaled_h_in
            scale_in_per_px *= shrink
            scaled_h_in = total_h_px * scale_in_per_px
        scaled_w_in = per_w_px * scale_in_per_px

        slide = prs.slides.add_slide(blank_layout)
        x_in = margin_in + (avail_w_in - scaled_w_in) / 2
        y_in = margin_in
        for png, im in zip(pngs, imgs):
            h_in = im.height * scale_in_per_px
            slide.shapes.add_picture(
                BytesIO(png),
                _emu(x_in),
                _emu(y_in),
                width=_emu(scaled_w_in),
                height=_emu(h_in),
            )
            y_in += h_in
        made += 1

    if made == 0:
        return 0
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path.as_posix())
    return made
